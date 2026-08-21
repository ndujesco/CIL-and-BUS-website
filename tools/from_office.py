"""PowerPoint and Word → blocks."""

import re

from docx import Document
from docx.table import Table as DocxTable
from docx.text.paragraph import Paragraph
from pptx import Presentation
from pptx.util import Emu

from common import Run, esc, runs_html, tidy


# ── PowerPoint ────────────────────────────────────────────────────────
def _runs(para):
    out = []
    for r in para.runs:
        out.append(Run(r.text, bool(r.font.bold), bool(r.font.italic)))
    return out


def _shape_paras(shape):
    """Paragraphs of a text frame, each as (indent level, html)."""
    out = []
    for p in shape.text_frame.paragraphs:
        html = runs_html(_runs(p))
        html = re.sub(r"^(?:[●•▪◦‣✓✔]|[-–—](?=\s))\s*", "", html)
        if html:
            out.append((p.level or 0, html))
    return out


def _table_html(shape):
    rows = []
    for n, row in enumerate(shape.table.rows):
        cell = "th" if n == 0 else "td"
        rows.append("<tr>" + "".join(
            "<" + cell + ">" + esc(tidy(c.text)) + "</" + cell + ">"
            for c in row.cells) + "</tr>")
    return "<table>" + "".join(rows) + "</table>"


def deck(path):
    """One slide per slide: its title, then its body text in reading order."""
    pres = Presentation(path)
    out = []
    for n, slide in enumerate(pres.slides, 1):
        title, body = "", []
        shapes = sorted(slide.shapes,
                        key=lambda s: (s.top if s.top is not None else 0,
                                       s.left if s.left is not None else 0))
        for shape in shapes:
            if shape.has_table:
                body.append((0, _table_html(shape)))
                continue
            if not shape.has_text_frame or not shape.text_frame.text.strip():
                continue
            is_title = (shape == slide.shapes.title
                        or (not title and shape.top is not None
                            and shape.top < Emu(1100000)
                            and len(shape.text_frame.text) < 90))
            paras = _shape_paras(shape)
            if is_title and paras and not title:
                title = re.sub(r"</?b>", "", " ".join(h for _, h in paras))
                continue
            body += paras
        if title or body:
            out.append((n, title, body))
    return out


# ── Word ──────────────────────────────────────────────────────────────
def _docx_table(table):
    rows = []
    for n, row in enumerate(table.rows):
        cell = "th" if n == 0 else "td"
        cells = []
        for c in row.cells:
            txt = tidy("\n".join(p.text for p in c.paragraphs))
            cells.append("<" + cell + ">" + esc(txt).replace("\n", "<br>")
                         + "</" + cell + ">")
        rows.append("<tr>" + "".join(cells) + "</tr>")
    return "<table>" + "".join(rows) + "</table>"


def _numbered(p):
    return p._p.pPr is not None and p._p.pPr.numPr is not None


def notes(path, masthead=0):
    """A Word document with no heading styles: bold short lines are the headings.

    The first few paragraphs are the department masthead, which is lifted into
    the document's metadata rather than repeated in its body.
    """
    doc = Document(path)
    meta, out = [], []
    for child in doc.element.body.iterchildren():
        if child.tag.endswith("}tbl"):
            out.append(("table", _docx_table(DocxTable(child, doc))))
            continue
        if not child.tag.endswith("}p"):
            continue
        p = Paragraph(child, doc)
        txt = tidy(p.text)
        if not txt:
            continue
        if len(meta) < masthead:
            meta.append(txt)
            continue

        html = runs_html(_runs(p))
        if _numbered(p):
            out.append(("oli" if re.match(r"^\d", txt) else "li", html))
            continue
        style = (p.style.name or "").lower()
        bold = all(r.bold for r in p.runs if r.text.strip())
        if style.startswith("heading 1") or (bold and re.match(r"^\d+\.\s", txt)):
            out.append(("h2", html))
        elif style.startswith("heading") or (bold and len(txt) < 80
                                             and not txt.endswith(".")):
            out.append(("h3", html))
        else:
            out.append(("p", html))
    return meta, out
